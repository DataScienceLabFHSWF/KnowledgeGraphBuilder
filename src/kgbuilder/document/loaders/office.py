"""Office document loaders (DOCX, PPTX).

Implementation of Issue #2.3: Office Document Loaders

TODO (DOCX):
- [x] Extract text from paragraphs
- [ ] Preserve document structure (headings, lists, bold/italic)
- [ ] Extract tables as structured data
- [ ] Handle embedded images and shapes
- [ ] Extract headers/footers
- [ ] Add unit tests

TODO (PPTX):
- [x] Extract text from slides
- [ ] Extract speaker notes
- [ ] Preserve slide order and hierarchy
- [ ] Extract embedded images
- [ ] Handle slide transitions/timing (optional)
- [ ] Add unit tests

See Planning/ISSUES_BACKLOG.md Issue #2.3 for acceptance criteria.
"""

from __future__ import annotations

from pathlib import Path
from uuid import uuid4

from kgbuilder.core.exceptions import DocumentLoadError
from kgbuilder.core.models import Document, DocumentMetadata, FileType


class DOCXLoader:
    """DOCX document loader using python-docx."""

    @property
    def supported_extensions(self) -> list[str]:
        """Return supported extensions."""
        return [".docx", ".DOCX"]

    def load(self, file_path: Path) -> Document:
        """Load DOCX document.

        Args:
            file_path: Path to DOCX file

        Returns:
            Document with extracted text

        Raises:
            DocumentLoadError: If DOCX cannot be read
        """
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise DocumentLoadError(
                str(file_path), "python-docx not installed"
            ) from None

        try:
            doc = DocxDocument(file_path)
            text_content = [para.text for para in doc.paragraphs if para.text]
            content = "\n\n".join(text_content)

            return Document(
                id=str(uuid4()),
                content=content,
                source_path=file_path,
                file_type=FileType.DOCX,
                metadata=DocumentMetadata(title=file_path.stem),
            )
        except Exception as e:
            raise DocumentLoadError(str(file_path), str(e)) from e

    def load_batch(self, file_paths: list[Path]) -> list[Document]:
        """Load multiple DOCX files."""
        documents = []
        for file_path in file_paths:
            try:
                doc = self.load(file_path)
                documents.append(doc)
            except DocumentLoadError:
                continue
        return documents


class PPTXLoader:
    """PPTX document loader using python-pptx."""

    @property
    def supported_extensions(self) -> list[str]:
        """Return supported extensions."""
        return [".pptx", ".PPTX"]

    def load(self, file_path: Path) -> Document:
        """Load PPTX document.

        Args:
            file_path: Path to PPTX file

        Returns:
            Document with extracted text

        Raises:
            DocumentLoadError: If PPTX cannot be read
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise DocumentLoadError(
                str(file_path), "python-pptx not installed"
            ) from None

        try:
            prs = Presentation(file_path)
            text_content = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content.append(shape.text)

            content = "\n\n".join(text_content)

            return Document(
                id=str(uuid4()),
                content=content,
                source_path=file_path,
                file_type=FileType.PPTX,
                metadata=DocumentMetadata(title=file_path.stem),
            )
        except Exception as e:
            raise DocumentLoadError(str(file_path), str(e)) from e

    def load_batch(self, file_paths: list[Path]) -> list[Document]:
        """Load multiple PPTX files."""
        documents = []
        for file_path in file_paths:
            try:
                doc = self.load(file_path)
                documents.append(doc)
            except DocumentLoadError:
                continue
        return documents
